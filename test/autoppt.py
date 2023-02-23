from pptx import Presentation
from pptx.util import Inches
import os

ppt=Presentation()

layout=ppt.slide_layouts[6]  #空白布局
slide=ppt.slides.add_slide(layout)

#文本框位置和大小
left=Inches(0)
top=Inches(0)
width=Inches(10)
height=Inches(7.5)# 4:3

textbox=slide.shapes.add_textbox(left,top,width,height)
textbox.text='新文本框'
textbox=slide.shapes.add_textbox(left,top,width,height)
textbox.text='文本框'

# new_paragrph=textbox.text_frame.add_paragraph()
# new_paragrph.text='文本框第二段'


name='插入文本框.pptx'
if os.path.exists(name):
    os.remove(name)
ppt.save(name)


        