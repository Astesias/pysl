# -*- coding: utf-8 -*-
"""
Created on Mon Apr 25 20:00:26 2022

@author: ysl

useful functions ysl created
"""

from __future__ import division
from __future__ import print_function
from __future__ import absolute_import
import types

class LazyLoader(types.ModuleType):

    """Lazily import a module, mainly to avoid pulling in large dependencies.

    `contrib`, and `ffmpeg` are examples of modules that are large and not always
    needed, and this allows them to only be loaded when they are used.
    """

    # The lint error here is incorrect.
    def __init__(self, local_name, parent_module_globals, name):  # pylint: disable=super-on-old-class

        self._local_name = local_name
        self._parent_module_globals = parent_module_globals

        super(LazyLoader, self).__init__(name)

    def _load(self):
        import importlib
        # Import the target module and insert it into the parent's namespace
        module = importlib.import_module(self.__name__)
        self._parent_module_globals[self._local_name] = module

        # Update this object's dict so that if someone keeps a reference to the
        #   LazyLoader, lookups are efficient (__getattr__ is only called on lookups
        #   that fail).
        self.__dict__.update(module.__dict__)

        return module

    def __getattr__(self, item):
        module = self._load()
        return getattr(module, item)

    def __dir__(self):
        module = self._load()
        return dir(module)

def lazyimport(module, asnmae=None):
    if not asnmae:
        asnmae = module
    return LazyLoader(asnmae, globals(), module)
    # from pysl import lazyimport
    # ti=lazyimport('time','ti')

# LazyLoader("chatglm_6B_int4", globals(), "IO_output.text.nlp.chatglm_6B_int4")
# lazyimport('torch')

import os
import re
import cv2
import sys
import json
import time

import shutil
import random
import logging
import threading
import traceback
import numpy as np
from functools import wraps
from contextlib import contextmanager

import requests
import websocket
from urllib.parse import quote
from bs4 import BeautifulSoup as bs
from requests_toolbelt import MultipartEncoder



#import pysl
# from pysl import *
# from pysl import My_list as ml
# from pysl import mmap

# String
# log_dir=r'D:\Desktop\.py\log'
# Chinese_Sign=r"[(|)|\n|-|！|，|。|？|、|（|）|￥|……|【|】| ]"
# colors={'cyan':(0,255,255),'blue':(0,0,255),'white':(255,255,255),
# 'red':(255,0,0),'yellow':(255,255,0),'green':(0,255,0),'black':(0,0,0)}  rgb->bgr

# matplotlib set
# plt.rc('savefig',dpi=500)
# plt.rc('legend',fontsize='large')
# plt.rc('axes',titlesize='large')
# plt.rc('ytick',labelsize='large')
# plt.rcParams['font.sans-serif']=['SimHei']
# plt.rcParams['axes.unicode_minus']=False

# pandas set
# pd.set_option('display.max_columns', None)
# pd.set_option('display.max_rows', None)
# pd.set_option('display.unicode.ambiguous_as_wide', True)
# pd.set_option('display.unicode.east_asian_width', True)

# Classes:
# excel_access                    #excel 操作
# Fplog                           #file io log
# ezlog                           #logger
# trace_logger                    #trace logger
# for_list()                      #循环列表
# My_list                         #运算符重载列表
# range_precent()                 #进度条类
# MyThread                        #自定义线程类
# Timer                           # 定时器

# Functions:
# bgr2rgb                         #颜色类型转化
# battery                         #电池信息
# c_b                             #condition debug
# cmd                             #cmd shell
# cv2_imread(filePath)            #非正确后缀名读取图片
# counts(inputs)                  #字符统计
# cfmt_str                        #颜色字符串
# cfmt_print                      #颜色格式输出
# cut_video                       #视频流转图片
# D(x)                            #方差
# drawRect                        #画矩形
# drewlinecross                   #画十字
# draw_points                     #画点
# del_list_by_index               #多索引删除
# easytxt                         #txt io操作
# easy_show_img                   #显示图片
# easy_request                    #url requests
# get_fname                       #当前文件
# getime                          #返回时间
# getip()                         #获取本机ip
# get_lll(alist)                  #二维列表转置
# get_son                         #文件递归获取
# isChinese(chs)                  #中文检测
# image_writer                    #图片下载
# join_list(l,sep=' '):           #列表转字符串
# list_index_in_list              #列表元素序列匹配
# min_index(alist)                #多个最小值索引
# mins(alist,n)                   #多个最小值
# mvdir(fm,to,showdetail=0):      #移动文件夹
# mmap                            #拓展map函数
# newpath                         #打开新文件夹文件
# path2filename(path)             #路径转文件名
# path2dirname(path)              #路径转目录名
# filename2end()                  #字符串转后缀
# printer                         #智能输出
# pic_insert_pic                  #图片插入
# pic_turn                        #图片处理
# qrcode                          #二维码
# rename123                       #批量重命名
# rotate_bound(image, angle)      #图片旋转
# random_color                    #随机颜色
# random_name                     #随机变量名
# re_args_get                     #args get
# re_findall                      #extend re
# ResizeImage                     #resize
# sorts                           #多条件排序
# say                             #语音播报
# sum_lll(alist)                  #列表列求和
# transparent_png                 #图片透明化
# tensor2img                      #Tensor图片显示
# truepath                        #模块中使用相对路径
# write_file_name                 #目录下文件名写入
# zipath(zipfile_path,Date_file)  #目录文件压缩
# zipsort                         #对应列表组合排序

# Decorations
# timeit                          #行数时计
# timety_                         #函数时计
# LazyLoader                      #延迟导入

########################################################################################################


class Auto_model():
    def __init__(self, model_dir, model_name_like, model_extend_name='.pth'):
        self.model_dir = model_dir
        self.model_name_like = model_name_like
        self.model_extend_name = model_extend_name
        newpath(os.path.join(model_dir, model_name_like))

    def create_new(self, n=0):
        return os.path.join(self.model_dir, self.model_name_like+f'{n}{self.model_extend_name}')

    def auto_load(self, get_latest=False):
        all_model = os.listdir(self.model_dir)
        latest = None

        for _ in all_model:

            n = _.strip(self.model_name_like).strip(self.model_extend_name)
            n = int(n)
            if not latest or n > latest:
                latest = n
        if get_latest:
            return latest
        if latest != None:
            r = os.path.join(self.model_dir, self.model_name_like +
                             f'{n}{self.model_extend_name}')
            print(f'load model {r}')
            return r
        else:
            print('no model exist')
            return None

    def auto_save(self, n):
        r = os.path.join(self.model_dir, self.model_name_like +
                         f'{n}{self.model_extend_name}')
        return r


class backref_dict(dict):
    def __init__(self, d, k2v=False):
        if k2v:
            assert len(list(d.values())) == len(
                set(d.values())), 'No element can duplication'
        self.k2v = k2v
        super(backref_dict, self).__init__(d)

    def backref(self, vaule):
        re = []
        for i, j in self.items():
            if j == vaule and self.k2v:
                return i
            elif j == vaule:
                re.append(i)
        return re

class Config():
    def __init__(self, path='./config_temp.json', base=None):
        self.path = path
        self.base = base
        if os.path.exists(path):
            with open(path) as fp:
                self.data = json.load(fp)
            for k, v in self.data.items():
                self.__setattr__(k, v)
        else:
            self.data = {}

    def load(self):
        if os.path.exists(self.path):
            with open(self.path) as fp:
                self.data = json.load(fp)
            for k, v in self.data.items():
                self.__setattr__(k, v)
    
    def add_base(self):
        for k, v in self.base.items():
            self.data[k]=v

    def add(self, item_k, item_v=None):
        assert self.base or item_v
        if item_v!=None:
            self.data[item_k] = item_v
            self.__setattr__(item_k, item_v)
        else:
            self.data[item_k] = self.base[item_k]
            self.__setattr__(item_k, self.base[item_k])

    def clean(self):
        self.data = {}

    def save(self, path=None):
        if not path:
            path = self.path
        with open(path, 'w') as fp:
            json.dump(self.data, fp, ensure_ascii=False, indent=2)

    def __getitem__(self,ind):
        return self(ind)

    def __enter__(self):
        return self

    def __exit__(self, type, value, trace):
        self.save()

    def __call__(self, key):
        return self.data[key]

    def __del__(self):
        try:
            self.save()
        except:
            pass


class dir_enter:
    def __init__(self, dir_=None):
        if not dir_:
            dir_ = truepath(__file__)
        self.dir = dir_

    def __enter__(self):
        return self

    def __exit__(self, type, value, trace):
        pass

    def _lisdir(self):
        return os.listdir(self.dir)

    def _realpath(self, path):
        return os.path.join(self.dir, path)

    def file(self, name):
        return os.path.join(self.dir, name)

    def md(self, path):
        os.mkdir(self._realpath(path))

    def cp(self, old, new):
        shutil.copy(self._realpath(old), self._realpath(new))

    def mv(self, old, new):
        shutil.move(self._realpath(old), self._realpath(new))

    def rm(self, path):
        if os.path.isfile(self._realpath(path)):
            os.remove(self._realpath(path))
        else:
            shutil.rmtree(self._realpath(path))

    def cd(self, path):
        self.dir = os.path.abspath(self._realpath(path))

    @property
    def ls(self):
        for _ in self._lisdir():
            __ = os.path.join(self.dir, _)
            if os.path.isfile(__):
                cfmt_print(f'@e@{_}', end=' ')
            else:
                cfmt_print(f'@^B@{_}', end=' ')
        print()

class os_enter:
    def __init__(self, dir_):
        dir_ = truepath(dir_)
        self.dir=self.odir= dir_

    def __enter__(self):
        return self

    def __exit__(self, type, value, trace):
        os.chdir(self.odir)

    def _lisdir(self):
        return os.listdir(self.dir)

    def file(self, name):
        return os.path.join(self.dir, name)

    def md(self, path):
        try:
            os.mkdir(path)
        except:
            print(os.path.abspath(path),'is existing')

    def cp(self, old, new):
        shutil.copy(old,new)

    def mv(self, old, new):
        shutil.move(old, new)

    def rm(self, path):
        if os.path.isfile(path):
            os.remove(path)
        else:
            shutil.rmtree(path)

    def cd(self, path):
        self.dir = os.chdir(path)

    @property
    def ls(self):
        for _ in self._lisdir():
            if os.path.isfile(_):
                cfmt_print(f'@e@{_}', end=' ')
            else:
                cfmt_print(f'@^B@{_}', end=' ')
        print()
    
    def cmd(self,order):
        os.system(order)


class ezlog():
    def __init__(self, filename):
        logging.basicConfig(
            filename=filename,
            # %(module)s.py line %(lineno)d
            format='%(asctime)s : %(message)s',
            datefmt='%Y/%m/%d|%H:%M:%S',
            level=logging.INFO
        )
        self.filename = filename

    def add(self, message):
        logging.critical(message)

    def err(self):  # useless in import module
        logging.error(traceback.print_exc())

    def flush(self):
        f = open(self.filename, 'w')
        f.close()


class Fplog():
    def __init__(self, filename):
        self.filename = filename+'_'+getime()+'.txt'
        newpath(self.filename)
        self.fp = open(self.filename, 'w')

    def add(self, message, t=True):
        self.fp = open(self.filename, 'a+')
        self.fp.write('{}{}\n'.format(getime()+' : ' if t else '', message))
        self.fp.close()

    def close(self):
        self.fp.close()

    def __del__(self):
        self.close()


class excel_access:
    # TODO
    def __init__(self, path, mode, sheet_num=0, begin_row=0, begin_line=0):  # 起始x y 0 0 全选
        import xlrd
        from openpyxl import load_workbook
        self.path = path
        self.mode = mode
        self.begin_row = begin_row
        self.begin_line = begin_line
        if mode == 'r':
            self.work_book = xlrd.open_workbook(path)
            self.table = self.work_book.sheets()[sheet_num]
        elif mode == 'w' or mode == 'n':
            self.work_book = load_workbook(path)
            sheetnames = self.work_book.get_sheet_names()  # 获得表单名字
            self.table = self.work_book.get_sheet_by_name(sheetnames[0])

    def get_cell_value(self, x, y):                               # x,y大于0
        try:
            return self.table.cell_value(y-self.begin_line-1, x-self.begin_row-1)
        except IndexError:
            print('indexerr in excel')
            return ''

    def get_cell_num(self, x, y):                               # x,y大于0
        try:
            return int(self.table.cell_value(y-self.begin_line-1, x-self.begin_row-1))
        except IndexError:
            print('indexerr in excel')
            return ''

    def write_cell(self, x, y, obj):
        from openpyxl.styles import Alignment
        self.table.cell(y-self.begin_line, x-self.begin_row,
                        obj).alignment = Alignment(horizontal='center', vertical='center')

        try:
            if self.mode == 'w':
                self.work_book.save(path2filename(self.path))
            elif self.mode == 'n':
                new_name = 'new_'+path2filename(self.path)
                self.work_book.save(new_name)
        except:
            raise IOError('文件被占用或已存在')


class for_list(list):
    def __init__(self, l):
        super(for_list, self).__init__(l)
        self.list = l

    def __getitem__(self, index):
        return self.list[index % len(self.list)]


class My_list(list):
    def __init__(self, alist):
        super(list, self).__init__()
        self.alist = list(alist)
        self.__radd__ = self.__add__
        self.__rmul__ = self.__mul__
        self.__rsub__ = self.__sub__
        self.__rtruediv__ = self.__truediv__

    def __add__(self, num):
        if type(self.alist) == type(self):
            self.list = self.alist.relist
        for n, i in enumerate(self.alist):
            self.alist[n] += num
        return self.alist

    def __mul__(self, num):
        if type(self.alist) == type(self):
            self.list = self.alist.relist
        if type(self.alist[0]) == type(''):
            l = self.alist[:]
            for i in range(num-1):
                self.alist.extend(l)
        else:
            for n, i in enumerate(self.alist):
                self.alist[n] *= num
        return self.alist

    def __sub__(self, num):
        if type(self.alist) == type(self):
            self.list = self.alist.relist
        for n, i in enumerate(self.alist):
            self.alist[n] -= num
        return self.alist

    def __truediv__(self, num):
        if type(self.alist) == type(self):
            self.list = self.alist.relist
        for n, i in enumerate(self.alist):
            self.alist[n] /= num
        return self.alist

    def relist(self):
        return self.alist


class range_percent():
    def __init__(self, total, process_name='Process', obj="█", nonobj='░', ef=True):
        self.total = total
        self.process_name = process_name
        self.obj = obj
        self.nonobj = nonobj
        self.ef = ef

    def update(self, now, new=''):
        precent = now/self.total
        num = int(100*precent)
        sys.stdout.flush()
        print("\r\r\r", end="")
        print("{} {:>3}% |".format(self.process_name, num), self.obj*(num//3),
              self.nonobj*(33-num//3), '|{}/{}'.format(now, self.total), sep='', end=new)
        if now == self.total and self.ef:
            print()
            self.ef = 0
        sys.stdout.flush()


class MyThread (threading.Thread):
    def __init__(self, func_method,
                 instance=None, arg=[], kw={},
                 lock=False, counter=None, name=None, threadID=None):
        threading.Thread.__init__(self)
        self.instance = instance
        self.func = func_method
        self.lock = lock
        self.arg = arg
        self.kw = kw
        self.threadID = threadID
        self.name = name
        self.counter = counter
        self.threadLock = threading.Lock()
        self.result = None

    def run(self):

        # print ("开启线程： " + self.name)
        if self.lock:
            self.threadLock.acquire()
        if self.counter:
            time.sleep(self.counter)

        global result
        if type(self.func) == type(''):
            print(f'Thread {self.threadID} runing')
            result = eval(
                ('self.instance.{}(*{},**{})').format(self.func, self.arg, self.kw))

        else:
            result = self.func(*self.arg, **self.kw)

        if self.lock:
            self.threadLock.release()
        self.result = result

        # if self.instance:
        #     return self.instance


class trace_logger:
    def __init__(self, filename='log.txt',
                 level=logging.DEBUG,
                 format='%(asctime)s - %(levelname)s - %(message)s'):

        self.filename = filename
        logging.basicConfig(filename=filename, level=level, format=format)

    def exec(self):

        errorFile = open(self.filename, 'a')
        errorFile.write('\n{} log from pysl module\n'.format(getime()))
        errorFile.write(traceback.format_exc())
        errorFile.close()

    def flush(self):
        errorFile = open(self.filename, 'w')
        errorFile.close()


class timeit():
    def __init__(self, linesname=''):
        self.info = linesname+':'
        self.line = sys._getframe().f_back.f_lineno+1
        self.t = time.time()

    def out(self, newinfo=None):
        line = sys._getframe().f_back.f_lineno
        print('In file {}| {} line {}-{} takes {} s'
              .format(get_fname(), self.info, self.line, line-1, time.time()-self.t))
        if newinfo:
            self.info = newinfo+':'
        self.t = time.time()
        self.line = line+1


class Timer():
    def __init__(self, sep=1):
        self.start = 0
        self.sep = sep

    def T(self):
        if time.time()-self.start > self.sep:
            self.start = time.time()
            return True
        else:
            return False
        
class Timer_12():
    def __init__(self, t1,t2):
        self.t1 = t1
        self.t2 = t2
        self.flag = False
        self.start = 0
        self.first_call = True

    def T(self):
        if self.first_call:
            self.first_call = False
            return True
        if time.time()-self.start < (self.t1 if self.flag else self.t2):
            return self.flag
        else:
            self.flag = not self.flag
            self.start = time.time()

class ppt_access():

    def __init__(self):
        from pptx import Presentation

        self.ppt = Presentation()
        self.slide_num = 0
        self.wid = 10
        self.hei = 7.5

    def add_slide(self, layout=6):
        self.slide_num += 1
        layout = self.ppt.slide_layouts[layout]
        slide = self.ppt.slides.add_slide(layout)
        self.slide = slide

    def add_textbox(self, text, left, top, width=5, height=0.25, ft=10):
        from pptx.util import Inches, Pt
        assert 0 < left < self.wid, 0 < top < self.hei
        if left+width >= self.wid:
            width = self.wid-left

        left = Inches(left)
        top = Inches(top)
        width = Inches(width)
        height_ = Inches(height)  # 4:3
        textbox = self.slide.shapes.add_textbox(left, top, width, height_)

        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = text
        font = para.font
        font.name = '微软雅黑'    # 字体类型
        font.size = Pt(ft if ft < 25 else 25)    # 大小
        font.underline = False    # 下划线

    def save(self, name='auto_ppt.pptx'):
        if os.path.exists(name):
            os.remove(name)
        self.ppt.save(name)


######################################################################################################

def admin_monitor(file):
    import ctypes
    ctypes.windll.shell32.ShellExecuteW(
        None, "runas", sys.executable, file, None, 1)


def args_sys():
    import sys
    return sys.argv[1:]


def args_dir2files(args):
    r = []
    for i in args:
        if os.path.isdir(i):
            for j in os.listdir(i):
                r.append(os.path.join(i, j))
        elif os.path.exists(i):
            r.append(i)
    return r

# TODO add_arg_bool sys.argv


def add_argments(dicts, help='argments get'):
    import argparse
    parser = argparse.ArgumentParser(description=help)
    for k, v in dicts.items():
        argname, types, *_ = v
        if _:
            default = _[0]
            if len(_) >= 2:
                help_ = _[-1]
            else:
                help_ = 'None'
        else:
            default = None
            help_ = 'None'

        parser.add_argument(argname, type=types, default=default, help=help_)
    return parser.parse_args()


def bgr2rgb(t):

    t = list(t)
    t[0], t[2] = t[2], t[0]
    return tuple(t)


def battery():
    import psutil
    a, b, c = psutil.sensors_battery()
    # print(a,b,c)
    print(f'当前电量: {a}% 预计可用 {b//60}分钟 ')

# TODO
# def colorComponents(hexAsString):
#     def uniformLength(strings):
#         return strings[:2] + "".zfill(10-len(strings)) + strings[2:]
#     re = []
#     re.append(str(hex(int(hexAsString, 16) & int("0xFF000000", 16))))
#     re.append(str(hex(int(hexAsString, 16) & int("0x00FF0000", 16))))
#     re.append(str(hex(int(hexAsString, 16) & int("0x0000FF00", 16))))
#     re.append(str(hex(int(hexAsString, 16) & int("0x000000FF", 16))))

#     for i in range(len(re)):
#         re[i] = uniformLength(re[i])
#     print(re)
#     return list(map(hex,re[1:]))


def clipboard(text):
    import win32clipboard as wcb
    import win32con as wc
    wcb.OpenClipboard()
    wcb.EmptyClipboard()
    wcb.SetClipboardData(wc.CF_TEXT, text.encode("utf8"))
    wcb.CloseClipboard()


def c_b(c):
    if c:
        print(0/0)


def cmd(command):
    import subprocess
    cmd = list(subprocess.getstatusoutput(command))
    # print(('Success' if not cmd[0] else 'Fail') + ' Command:\n   '+command)
    cmd[1]=cmd[1].replace('Active code page: 65001', '')
    # print(cmd)
    if cmd[0]==0:
        return cmd[1]
    else:
        print(f'Command \'{command}\' failed')



def cfmt_str(s):

    # f=[0,1,4,7] #  ^_~
    # c=[30,31,32,33,34,35,36,37] # black r g y blue p grey w    rgybpew

    fmtchar = list(' ^_~hrgybpewHRGYBPEW')
    mapchar = ['0', '1', '4', '7', '30', '31', '32', '33', '34', '35',
               '36', '37', '40', '41', '42', '43', '44', '45', '46', '47']
    assert (0 < len(s) <= 3 and s.split)
    s = list(s)
    for n, i in enumerate(s):
        assert i in fmtchar
        s[n] = mapchar[fmtchar.index(i)]

    ctrl = '\033[{}m'
    return ctrl.format(';'.join(s))


def cfmt_print(*fmtstrs, syschar='@', endchar='$', fmt=True, endfmt=True, return_=False, **kw):
    '''
     ^_~ 高亮 下划 反显
     hrgybpew 黑红绿黄靛紫蓝白
     HRGYBPEW
    '''
    clear = '\033[0m'
    fmtstrs = list(fmtstrs)
    for n, fmtstr in enumerate(fmtstrs):
        if endfmt:
            fmtstrs[n] += endchar
        fmtstrs[n] = fmtstrs[n].replace(endchar, clear)
        fmtstrs[n] = fmtstrs[n].replace(syschar, '@')
        r = re.findall(r'@.{1,3}@', fmtstrs[n])
        # print(r)
        for i in r:
            j = i.strip(syschar)
            fmtstrs[n] = fmtstrs[n].replace(i, cfmt_str(j))
        if fmt:
            r = re.findall(r'{.+?}', fmtstrs[n])
            for i in r:
                j = i[1:][:-1]
                fmtstrs[n] = fmtstrs[n].replace(i, str(eval(j)))

    if not return_:
        print(*fmtstrs, **kw)
    else:
        return ''.join(fmtstrs)

    # cfmt_print('website ：@_bE@{head}.@r~p@github{111*6}.@^hG@com')


def cv2_imread(filePath):
    cv_img = cv2.imdecode(np.fromfile(filePath, dtype=np.uint8), -1)
    return cv_img


def counts(inputs, sep=' '):
    a = inputs
    d = {}
    b = a.split(sep)
    for i in b:
        try:
            d[i] += 1
        except:
            d[i] = 1

    s = list(d.items())
    s.sort(key=lambda x: x[1], reverse=True)

    return s


def cut_video(i_video, o_video, skip, expand_name='.jpg'):
    cap = cv2.VideoCapture(i_video)
    num_frame = cap.get(cv2.CAP_PROP_FRAME_COUNT)
    print('expect', int(num_frame//skip), 'output picture')
    if not cap.isOpened():
        print("Please check the path.")
    cnt = 0
    count = 0
    while 1:
        ret, frame = cap.read()
        cnt += 1
        if cnt % skip == 0:
            count += 1
            # print(o_video, str(count) + expand_name,num_frame)
            try:
                cv2.imwrite(os.path.join(
                    o_video, str(count) + expand_name), frame)
            except:
                pass
        if not ret:
            break
    print('cut_video done')

class response_wrapper():
    def __init__(self,response):
        self.r = response
    def response(self):
        return self.r
    def byte(self):
        return self.r.content
    def str(self,encoding='utf8'):
        return str(self.byte(), encoding=encoding)
    def json(self,*args,**kws):
        return json.loads(self.str(*args,**kws))
    def bs(self,*args,**kws):
        return bs(self.str(*args,**kws), features='lxml')
    def byte_iter(self,chunk_size=1024):
        return self.r.iter_content(chunk_size=chunk_size)
    def __repr__(self):
        return self.r.__repr__()
    

class chained_request():
    def __init__(self,url):
        self._url = url
        self._headers = None
        self._payload = None
        self._urlencode = False
        
    def quote(self,format_url_args):
        self._url = self._url.format(*list(map(quote, format_url_args)))
        return self
    
    def payload(self,data):
        self._payload = data
        return self
    
    def urlencode(self):
        self._urlencode = True
        return self
    
    def mutipayload(self,fields,boundary):
        custom_data = MultipartEncoder(
            fields=fields, boundary=boundary)
        if not self._headers:
            self._headers = {}
        self._headers['Content-Type'] = custom_data.content_type
        self._payload = custom_data
        return self
    
    def get(self,stream=False):
        return self._request('GET',stream = stream)
    
    def post(self,payload=None,stream=False):
        if payload:
            self._payload = payload
        return self._request('POST',stream = stream)
    
    def headers(self,path_or_dict):
        if isinstance(path_or_dict,str):
            self._headers = json.load(open(path_or_dict))
        else:
            self._headers = path_or_dict
        return self
    
    def _request(self,method,stream=False):
        if method=='POST':
            response = requests.post(self._url,
                                     headers=self._headers,
                                     data=self._payload if self._urlencode else json.dumps(self._payload),
                                     stream = stream
                                     )
        elif method=='GET':
            response = requests.get(self._url,
                                     headers=self._headers,
                                     stream = stream
                                     )
        return response_wrapper(response)
    
class chain_ws:
    def __init__(self,url):
        self.url = url
        self.ws = websocket.create_connection(url)
        self._payload = None
    def payload(self,payload):
        self._payload = payload
        return self
    def recv(self):
        if r:=self.ws.recv():
            return r
        else:
            return
    def recv_all(self):
        res = []
        while 1:
            r = self.ws.recv()
            if r:
                res.append(r)
            else:
                break
        return [response_wrapper(_) for _ in res]
    def send(self,payload=None):
        if payload:
            self._payload = payload
        self.ws.send(json.dumps(self._payload))
        return self
    def close(self):
        self.ws.close()
    def __del__(self):
        self.close()
    
def colorhex2rgb(s):
    s = s.replace('#','')
    s1=s[:2]
    s2=s[2:4]
    s3=s[4:6]
    return int('0x'+s1,16),int('0x'+s2,16),int('0x'+s3,16)
    

def D(x):  # 方差
    d = 0
    avg = sum(x)/len(x)
    for i in range(len(x)):
        d += (x[i]-avg)**2
    return d


def dir_search(clas, t, uplow=True):
    for i in dir(clas):
        if not uplow:
            if t in i:
                print(i)
        else:
            if t.lower() in i.lower():
                print(i)


def dow(y, m, d):
    t = [0, 3, 2, 5, 0, 3, 5, 1, 6, 4, 2, 4]
    y -= m < 3
    return int((y+y/4-y/100+y/400+t[m-1]+d) % 7)


def drawtri(img, pt1, pt2, pt3, pt4, color=(255, 255, 255), lineWidth=2):
    # cv2.line(img, pt1, pt2, color, lineWidth)
    cv2.line(img, pt2, (pt2[0], pt3[1]), (0, 0, 255), lineWidth)
    cv2.line(img, pt2, pt3, (0, 0, 255), lineWidth)
    cv2.line(img, (pt2[0], pt3[1]), pt3, color, lineWidth)


def drawRect(img, pt1, pt2, pt3, pt4, color=(255, 255, 0), lineWidth=2):
    pt1 = tuple(map(int, pt1))
    pt2 = tuple(map(int, pt2))
    pt3 = tuple(map(int, pt3))
    pt4 = tuple(map(int, pt4))
    cv2.line(img, pt1, pt2, color, lineWidth)
    cv2.line(img, pt2, pt3, color, lineWidth)
    cv2.line(img, pt3, pt4, color, lineWidth)
    cv2.line(img, pt1, pt4, color, lineWidth)


def drewlinecross(img, x, mode='l', lineWidth=2, color=(255, 255, 0)):
    h, w, d = img.shape
    if mode == 'l':
        cv2.line(img, (x, 0), (x, h), color, lineWidth)
    else:
        cv2.line(img, (0, x), (w, x), color, lineWidth)


def draw_points(*ps, img='D:\\Desktop\\.py/block.jpg', resize=(900, 900)):
    img = cv2_imread(img)
    if resize:
        img = cv2.resize(img, resize)
    for n, p in enumerate(ps):
        # print(p)
        x = p[0]
        y = p[1]
        p1 = (x-2, y-2)
        p2 = (x+2, y-2)
        p3 = (x+2, y+2)
        p4 = (x-2, y+2)
        p1 = tuple(map(int, p1))
        p2 = tuple(map(int, p2))
        p3 = tuple(map(int, p3))
        p4 = tuple(map(int, p4))

        drawRect(img, p1, p2, p3, p4, color=bgr2rgb((255, 0, 0)), lineWidth=2)
        # print((p1[0]+6,p[1]-6))
        cv2.putText(img, 'p'+str(n+1), (p1[0]+6, p1[1]-6),
                    cv2.FONT_HERSHEY_TRIPLEX, 0.4, bgr2rgb((255, 0, 0)), 1)
    cv2.imshow('points', img)
    cv2.waitKey(0)
    cv2.destroyAllWindows()


def del_list_by_index(l, i):
    l2 = l.copy()
    for j in i:
        l.remove(l2[j])
    return l


def easytxt(data, txtname, mode='w'):
    if txtname[-4:] != '.txt':
        txtname += '.txt'
    if mode == 'w' or mode == 'x':
        with open(txtname, mode) as fp:
            if type(data) == type([]):
                for i in data:
                    fp.write(i)
                    fp.write('\n')
            else:
                fp.write(data)
    else:
        return fp.read()


def easy_show_img(img, rate=1, name=' ', transpose=None, bgr=False):
    if type(img) == type(''):
        img = cv2_imread(img)
    if transpose:
        img = img.transpose(*transpose)
    if bgr:
        try:
            b, g, r, _ = cv2.split(img)
            img = cv2.merge([r, g, b, _])
        except:
            b, g, r, = cv2.split(img)
            img = cv2.merge([r, g, b])
    if rate and rate != 1:
        h, w, d = img.shape
        img = cv2.resize(img, (int(rate*h), int(rate*w)))
    try:
        cv2.imshow(name, img)
    except:
        cv2.imshow(name, img.transpose(1, 2, 0))
    cv2.waitKey(0)
    cv2.destroyAllWindows()


def easy_request(url, header=None, format_url_args=None,
                 data=None, method='GET', driver=False, pic=False, form_data=None):

    from urllib.parse import quote
    import requests
    import json
    from bs4 import BeautifulSoup as bs

    if format_url_args:
        url = url.format(*list(map(quote, format_url_args)))
    if header and isinstance(header, str):
        with open(header) as fp:
            header = json.load(fp)
    if data:
        data = json.dumps(data)

    if not driver:
        if pic:
            return requests.request(method, url,
                                    headers=header,
                                    data=data,
                                    )
        if form_data:
            from requests_toolbelt import MultipartEncoder
            custom_data = MultipartEncoder(
                fields=form_data[0], boundary=form_data[1])
            if not header:
                header = {}
            header['Content-Type'] = custom_data.content_type
            response = requests.post(url,
                                     headers=header,
                                     data=custom_data,
                                     ).content
        else:
            response = requests.request(method, url,
                                        headers=header,
                                        data=data,
                                        ).content
    else:
        from selenium import webdriver
        browser = webdriver.Chrome()
        browser.get(url)
        data = browser.page_source
        browser.quit()
        soup = bs(data, features="lxml")
        response = soup.find('body').text
        return response

    try:
        if isinstance(response, bytes):
            response = str(response, encoding='utf8')
            return response
    except:
        pass
    try:
        data = json.loads(response)
        return data
    except:
        soup = bs(response, features='lxml')
        return soup


def flatten(l):
    r = []
    for i in l:
        r.extend(i)
    return r

def find_files_with_extension(directory, extension,with_root=False):
    result = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(extension):
                if with_root:
                    result.append(os.path.join(root, file))
                else:
                    result.append(file)
    return result

def getime(f=None):
    from datetime import datetime
    s = datetime.now().strftime('%Y-%m-%d-%H-%M-%S')
    p = s.split('-')
    if not f:
        return '{}-{}-{} {}:{}:{}'.format(*p)
    else:
        return '{}-{}-{}_{}{}{}'.format(*p)


def get_fname(ori=False):
    if ori:
        return path2filename(__file__)
    return path2filename(sys.argv[0])


def getip():
    import uuid
    import socket
    ip = socket.gethostbyname(socket.gethostname())
    node = uuid.getnode()
    macHex = uuid.UUID(int=node).hex[-12:]
    mac = []
    for i in range(len(macHex))[::2]:
        mac.append(macHex[i:i+1])

    mac = ':'.join(mac)
    return ip, mac


def get_lll(alist):  # 获取列
    l = len(alist[0])
    ll = len(alist)
    lllist = [[] for i in range(l)]
    for i in range(l):
        for j in range(ll):
            lllist[i].append(alist[j][i])

    return lllist


def get_son(path, get_all=False, judge=None, judge_mode='all', orgin_path=False, dir_choose=False):  # 获取子文件
    ''' 
    (路径，遍历所有文件标志，判断条件（字符串），不包括原路径的判断方式（开头，自己，结尾），
     展示原路径标志,展示文件夹路径标志（需开启get_all）)
    TODO judges
    '''
    if path[-1] == '/':
        path = path[:-1]
    if path[-2:] == '\\':
        path = path[:-2]

    list2 = []
    if get_all:
        for root, dirs, files in os.walk(path):
            list1 = []
            if orgin_path:
                for name in files:
                    list1.append(root+'\\'+name)
                    if dir_choose:
                        for dir in dirs:
                            list1.append(root+'\\'+dir)
                if judge:
                    if judge_mode == 'last':
                        for name in list1:
                            name_last2 = name[-len(judge):]
                            if judge == name_last2:
                                list2.append(name)
                    elif judge_mode == 'in':
                        for name in list1:
                            name_in = name[name.rfind('\\')+1:]
                            if judge in name_in:
                                list2.append(name)
                    elif judge_mode == 'begin':
                        for name in list1:
                            name_in = name[name.rfind('\\')+1:]
                            name_begin = name_in[:len(judge)]
                            if name_begin == judge:
                                list2.append(name)
                    else:
                        for name in list1:
                            name_in = name[name.rfind('\\')+1:]
                            if judge in name_in:
                                list2.append(name)
                else:
                    print(list1)
                    list2.extend(list1)

            else:
                for name in files:
                    list1.append(name)
                if dir_choose:
                    for dir in dirs:
                        list1.append(dir)
                if judge:
                    if judge_mode == 'last':
                        for name in list1:
                            name_last2 = name[-len(judge):]
                            if judge == name_last2:
                                list2.append(name)
                    elif judge_mode == 'in':
                        for name in list1:
                            name_in = name[1:-1]
                            if judge in name_in:
                                list2.append(name)
                    elif judge_mode == 'begin':
                        for name in list1:
                            name_begin = name[:len(judge)]
                            if judge == name_begin:
                                list2.append(name)
                    else:
                        for name in list1:
                            if judge in name:
                                list2.append(name)
                else:
                    list2.extend(list1)
        return list2

    else:
        list1 = os.listdir(path)
        li = list(list1)
        for i in li:
            if os.path.isdir(path+'\\'+i):
                list1.remove(i)
        if orgin_path:
            p = path + '\\'
        else:
            p = ''
    if judge:
        if judge_mode == 'last':
            for name in list1:
                name_last = name[-len(judge):]
                if judge == name_last:
                    list2.append(p+name)
        elif judge_mode == 'in':
            for name in list1:
                name_in = name[1:-1]
                if judge in name_in:
                    list2.append(p+name)
        elif judge_mode == 'begin':
            for name in list1:
                name_begin = name[:len(judge)]
                if judge == name_begin:
                    list2.append(p+name)
        else:
            for name in list1:
                if judge in name:
                    list2.append(p+name)
    else:
        for i, name in enumerate(list1):
            list1[i] = p+list1[i]
        return list1
    return list2


def get_file_md5(filename):
    import hashlib
    if not os.path.exists(filename):
        return False
    if not os.path.isfile(filename):
        return False
    md5 = hashlib.md5()
    f = open(filename, 'rb')
    while True:
        b = f.read(8096)
        if not b:
            break
        md5.update(b)
    f.close()
    return md5.hexdigest()


def get_file_same(file1, file2):
    a,b=get_file_md5(file1),get_file_md5(file2)
    if not (a and b):
        return False
    else:
        return a==b


def is_admin():
    import ctypes
    try:
        1/ctypes.windll.shell32.IsUserAnAdmin()
        return True
    except:
        return False


def isChinese(chs):
    if len(chs) == 1:
        return u'\u4e00' <= chs <= u'\u9fff'
    else:
        f = 0
        for ch in chs:
            if u'\u4e00' <= ch <= u'\u9fff':
                pass
            else:
                f += 1
        return f == 0


def img_writer(path='../../data/images', num=1, url=None, headers=None):
    names = []

    if not url:
        url = 'https://api.yimian.xyz/img'
    if isinstance(headers, str):
        with open(truepath(__file__, headers)) as fp:
            headers = json.load(fp)

    for i in range(num):
        name = random_name(5)+'.jpg'
        names.append(name)
        r = easy_request(url, header=headers, pic=True)
        p = truepath(__file__, path, name)
        _ = 1
        while _ or os.stat(p).st_size < 3000:
            if _:
                _ = 0
            else:
                time.sleep(1)
                print('.', end='')
                r = easy_request(url, header=headers, pic=True)
            with open(p, "wb") as f:
                f.write(r.content)
    return names


def join_list(l, sep=' '):
    s = ''
    for n, i in enumerate(l):
        if n != len(l)-1:
            s += str(i)
            s += sep
        else:
            s += str(i)
    return s


def list_index_in_list(al, bl):
    index = []
    for i in al:
        for n, j in enumerate(bl):
            if i == j:
                index.append(n)
                break
            if n == len(bl)-1:
                index.append(None)
    return index


def min_index(alist):  # 最小值索引列表
    indexes = []
    blist = list(alist)

    miner, indexm = min(blist), blist.index(min(blist))
    indexes.append(indexm)

    blist.remove(miner)
    i = 1
    while(miner in blist):
        miner, indexm = min(blist), blist.index(min(blist))
        indexes.append(indexm+i)
        i += 1
        blist.remove(miner)

    return indexes


def mins(alist, n):
    s = []
    for i in range(n):
        s.append(alist.index(min(alist)))
        alist[alist.index(min(alist))] = max(alist)+1
    return s


def mvdir(fm, to, showdetail=0):
    import shutil
    paths = os.listdir(fm)
    opath = [os.path.join(fm, i) for i in paths]
    tpath = [os.path.join(to, i) for i in paths]
    for n, (o, t) in enumerate(zip(opath, tpath)):
        shutil.copy(o, t)
        if showdetail:
            print('copy from {} to {}   total:{}'.format(o, t, n))
    print('复制 {} 文件夹至目标文件夹 {}'.format(fm, to))


def mmap(func_or_method, ite, arg=[], kw={}):
    r = []
    if type(arg) != type([]):
        arg = [arg]
    if type(func_or_method) == type(''):
        for i in ite:
            # print(('i.{}(*{},**{})'.format(func_or_method,arg,kw)))
            r.append(eval('i.{}(*{},**{})'.format(func_or_method, arg, kw)))
        return r
    else:
        for i in ite:
            r.append(func_or_method(i, *arg, **kw))
        return r

def make_dpi_aware():
    import platform,ctypes
    if int(platform.release()) >= 8:
        ctypes.windll.shcore.SetProcessDpiAwareness(True)

def make_header_json(name):
    import re
    with open('temp.json') as fp:
        data=fp.read()
    # data='\n'.join(data)
    print(data)
    
    keys=re.findall('.+?(?=:): ',data)
    values=re.findall('(?<= ).+',data)
    
    with open(name+'.json','w') as fp:
        fp.write('{\n')
        for i,j in zip(keys,values):
            i=i.strip(': ')
            fp.write(f'\"{i}\": \"{j}\",\n')
        
        fp.write(r'"py": "ysl"')
        fp.write('\n}')

def md5(input_string):
    import hashlib
    # 创建一个 MD5 hash 对象
    md5_hash = hashlib.md5()
    
    # 更新 hash 对象的输入，将输入字符串编码为字节型并更新到 hash 对象中
    md5_hash.update(input_string.encode('utf-8'))
    
    # 获取十六进制表示的散列值
    md5_digest = md5_hash.hexdigest()
    
    return md5_digest

@contextmanager
def mute_all():
    with open(os.devnull, "w") as devnull:
        old_stdout = sys.stdout
        sys.stdout = devnull
        try:
            yield
        finally:
            sys.stdout = old_stdout


def newpath(path, isfile=False):
    assert path[0] != '/'
    p = []
    for i in range(10):
        if not os.path.exists(path):
            l, r = os.path.split(path)
            p.append(r)
        else:
            break
        path = l[:]
    if isfile:
        p = p[1:]
    for _ in p[::-1]:
        l = os.path.join(l, _)
        os.mkdir(l)


def path2filename(path):
    if type(path) != type('str'):
        raise TypeError('path is a str,not {}'.format(type(path)))
    if path.rfind('\\') > path.rfind('/'):
        return path[path.rfind('\\')+1:]
    else:
        return path[path.rfind('/')+1:]


def filename2end(path, filp=0, find='.', include=True):
    if type(path) != type('str'):
        raise TypeError('path is a str,not {}'.format(type(path)))
    if include:
        if not filp:
            return path[path.rfind(find):]
        else:
            return path[:path.rfind(find)]
    else:
        if not filp:
            return path[path.rfind(find)+1:]
        else:
            return path[:path.rfind(find)+1]


def path2dirname(path):
    if type(path) != type('str'):
        raise TypeError('path is a str,not {}'.format(type(path)))
    a = path[:path.rfind('\\')+1]
    b = path[:path.rfind('/')+1]
    return a if len(a) > len(b) else b


def pic_insert_pic(
    back_pic, front_pic,
    output_path, output_form='path', data_form='path',  # str  图片大小位置随机插入
    insert_position=(0, 0),  # list
    resize=False, resize_mode=1, front_resize_size=None, back_resize_size=None,
    random_position=False, random_size=False,
    random_size_mini=None, random_size_max=None
):
    '''
    [
     背景(路径)，前景(路径)，
     输出路径，输出形式（文件/return），输出形式(路径/np)，
     插入位置（默认左上），前景变化标志，前景变化模式，前景变化大小，背景变化大小，
     随机位置标志，随机大小标志，
     随机最小值，随机最大值
     ]
    '''
    if data_form == 'path':
        pic_back = cv2_imread(back_pic)
        pic_front = cv2_imread(front_pic)
    else:
        pic_back = back_pic
        pic_front = front_pic
    if back_resize_size:  # 变化背景大小
        pic_back = cv2.resize(pic_back, back_resize_size,
                              interpolation=cv2.INTER_NEAREST)

    bx, by = pic_back.shape[0], pic_back.shape[1]
    fx, fy = pic_front.shape[0], pic_front.shape[1]

    if bx <= fx or by <= fy:
        print('Warning: pic_back is smaller than pic_front,resize would be nessnecessary ')
        if resize == False:
            raise IOError('Error: pic_back is smaller than pic_front')

    if random_size:  # 随机变化大小
        # if bx <= fx:
        #     fx=0.3*bx
        #     fy=0.3*by
        x2 = y2 = random.randint(random_size_mini[0], random_size_max[0])
        pic_front = cv2.resize(pic_front, (x2, y2),
                               interpolation=cv2.INTER_NEAREST)
        fx, fy = pic_front.shape[0], pic_front.shape[1]

    if resize and random_size == False:  # 按输入变化大小
        if resize_mode == 1:
            pic_front = cv2.resize(
                pic_front, front_resize_size, interpolation=cv2.INTER_NEAREST)
        elif resize_mode == 2:  # TODO
            pic_front = cv2.resize(
                pic_front, front_resize_size, interpolation=cv2.INTER_NEAREST)
        else:
            raise ValueError('Error: resize_mode should be 1 or 2')
        fx, fy = pic_front.shape[0], pic_front.shape[1]

    if random_position:  # 随机插入位置
        if insert_position != None:
            print(
                'Warning: in mode random_position,defined insert_position would not be used')
        x = random.randint(0, bx-fx)
        y = random.randint(0, by-fy)
        insert_position = (x, y)

    pic_back[insert_position[0]:insert_position[0]+fx,
             insert_position[1]:insert_position[1]+fy] = pic_front
    new_pic = pic_back
    if output_form == 'path':
        cv2.imwrite(output_path, new_pic)
    else:
        return new_pic


def pic_turn(path='C:/Users/ysl/Desktop/.py/pic.jpg', mode=0):

    pic = cv2_imread(path)
    if mode == 0:
        pic = cv2.flip(pic, 0)  # 倒影
    if mode == 1:
        pic = cv2.flip(pic, 1)  # 镜像
    if mode == 2:
        pic = cv2.flip(pic, -1)  # 倒影镜像
    if mode == 3:
        pic = cv2.cvtColor(pic, cv2.COLOR_RGB2BGR)  # RGB -> BGR
    if mode == 4:
        pic = cv2.cvtColor(pic, cv2.COLOR_RGB2GRAY)  # 黑白
    if mode == 5:
        # 模糊
        pass
    # 裁剪 拼接 旋转 降噪
    cv2.imshow('1', pic)
    cv2.waitKey(0)


def qrcode(word, pic=None, save="二维码.png"):
    from PIL import Image as img
    # word="https://www.baidu.com"
    # save="二维码.png"
    # pic=r"D:\Desktop\Image\6945ab7f7ef0153745264d990cc3274.png"

    # myqr.run(
    #         words = "https://www.baidu.com",
    #         level='H',
    #         picture = r"D:\Desktop\Image\6945ab7f7ef0153745264d990cc3274.png",
    #         colorized=True,
    #         save_name=save,
    #         )

    if pic:
        cmd(f'myqr {word} -v 4 -p {pic} -d {os.getcwd()} -n {save} -c')
    else:
        cmd(f'myqr {word} -v 4 -d {os.getcwd()} -n {save} ')

    i = img.open(save)
    i.resize((200, 200), img.ANTIALIAS)
    i.show()


def renames(path, new=None, nokey=None):
    #path =r'C:\Users\ysl\Desktop\新建文件夹'
    files = os.listdir(path)
    a = 0
    if nokey:
        pass
    else:
        nokey = '\\'

    oldn = []
    newn = []
    for i, file in enumerate(files):
        if nokey not in file:
            NewName = os.path.join(
                path, str(a)+(new if new else filename2end(file)))
            newn.append(NewName)
            OldName = os.path.join(path, file)
            oldn.append(OldName)
            os.rename(OldName, NewName)
            a += 1
    print('draw {} files'.format(a))
    for i in range(len(newn)):
        print(path2filename(oldn[i]), '->', path2filename(newn[i]))


def random_color():
    return (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))


def random_name(n):
    st = 'ABCDEFGHIJKLMNOPQRSTUVWXYZ'
    st = st+st.lower()+'_0123456789'
    r = ''
    for i in range(n):
        if i == 0:
            r += random.choice(st[:-9])
        else:
            r += random.choice(st)
    return r


def rotate_bound(image, angle):
    # grab the dimensions of the image and then determine the center
    # 抓取图像的尺寸，然后确定中心
    (h, w) = image.shape[:2]
    (cX, cY) = (w // 2, h // 2)
    # grab the rotation matrix (applying the negative of the angle to rotate clockwise), then grab the sine and cosine (i.e., the rotation components of the matrix)
    # 抓取旋转矩阵（应用角度的负数顺时针旋转），然后抓取正弦和余弦（即矩阵的旋转分量）
    M = cv2.getRotationMatrix2D((cX, cY), -angle, 1.0)
    cos = np.abs(M[0, 0])
    sin = np.abs(M[0, 1])
    # compute the new bounding dimensions of the image
    # 计算图像的新边界尺寸
    nW = int((h * sin) + (w * cos))
    nH = int((h * cos) + (w * sin))
    # adjust the rotation matrix to take into account translation
    # 调整旋转矩阵以考虑平移
    M[0, 2] += (nW / 2) - cX
    M[1, 2] += (nH / 2) - cY
    # perform the actual rotation and return the image
    # 执行实际旋转并返回图像
    return cv2.warpAffine(image, M, (nW, nH))


def re_matchall(pattern, s):
    r = []
    match = re.search(pattern, s)
    while match:
        r.append(match.group())
        s = s[0:match.span()[0]]+s[match.span()[1]:]
        match = re.search(pattern, s)
    return r


def re_args_get(msg, d, nolist=True):
    # dict{'arg_name':('as_name','type_name',default,is_single)}
    # types=['int','float_x','cnt','bool','str']
    result_d = {}
    for i in list(d.items()):
        arg_name = i[0]
        as_name = i[1][0]
        type_name = i[1][1]
        default = i[1][2]
        try:
            is_single = i[1][3]
        except:
            is_single = True
        assert nolist == is_single

        if type_name == 'cnt':
            match_ = msg.count(as_name)
            result_d[arg_name] = match_ if nolist else [match_]
            continue
        elif type_name == 'bool':
            match_ = not default if msg.count(as_name) else default
            result_d[arg_name] = match_ if nolist else [match_]
            continue
        elif type_name == 'bool+':
            match_ = default[1] if msg.count(as_name) else default[0]
            result_d[arg_name] = match_ if nolist else [match_]
            continue
        else:
            matches = re_matchall(f'(?<!\w){as_name}(=| )?.+?( |$)', msg)
            if len(matches) == 0:
                if type_name == 'bool+':
                    result_d[arg_name] = default[0] if nolist else [default[0]]
                else:
                    result_d[arg_name] = default if nolist else [default]

            elif is_single and len(matches) >= 2:
                raise Exception(f'single err in arg {as_name}')
            else:
                arg_list = []
                for match in matches:
                    msg = msg.replace(match, '')
                    match_ = match.replace(f'{as_name}=', '').\
                        replace(f'{as_name}', '').replace(f'{as_name} ', '')
                    if type_name == 'int':
                        match_ = int(match_)
                    elif type_name == 'str':
                        match_ = match_.strip(' ').strip('\n')
                    elif type_name.count('float'):
                        pcn = int(type_name[-1])
                        match_ = round(float(match_), pcn)
                    arg_list.append(match_)
                result_d[arg_name] = arg_list[0] if nolist else arg_list
    return result_d


def ResizeImage(filein, fileout, maxlimit=1024):
    from PIL import Image
    img = Image.open(filein)
    type_ = img.format

    width, height = img.size
    if width > height:
        height = maxlimit*height/width
        width = maxlimit
    else:
        width = maxlimit*width/height
        height = maxlimit

    out = img.resize((int(width), int(height)), Image.ANTIALIAS)
    out.save(fileout, type_)


def shorten_url(url):
    from pyshorteners import Shortener
    short_engine = Shortener()
    res = short_engine.tinyurl.short(url)
    print(res)


def sorts(*args, l, mode='n'):
    l1 = list(l)
    for i in args:
        l1.sort(key=lambda x: x[i])
    if mode == 'n':  # new
        return l1
    elif mode == 'c':  # cover
        l = l1
        return l


def say(sen, times=1, save=False):
    import pyttsx3
    engine = pyttsx3.init()  # 初始化
    engine.setProperty('rate', 200)  # 设置语速
    engine.setProperty('volume', 3.0)  # 设置音量
    voices = engine.getProperty('voices')
    engine.setProperty('voice', voices[0].id)
    # sen='傻逼'
    for i in range(times):
        engine.say(sen)
    if save:
        engine.save_to_file(sen, f"./{sen}.mp3")
    engine.runAndWait()


def screenxy(wl=0, hl=0):
    import subprocess
    try:
        cmd = ['xrandr']
        cmd2 = ['grep', '*']
        p = subprocess.Popen(cmd, stdout=subprocess.PIPE)
        p2 = subprocess.Popen(cmd2, stdin=p.stdout, stdout=subprocess.PIPE)
        p.stdout.close()
        resolution_string, junk = p2.communicate()
        resolution = resolution_string.split()[0]
        width, height = resolution.split(b'x')
        print('linux')
        return int(width)-wl, int(height)-hl
    except:
        from win32 import win32api, win32gui, win32print
        from win32.lib import win32con

        """获取缩放后的分辨率"""
        sX = win32api.GetSystemMetrics(0)  # 获得屏幕分辨率X轴
        sY = win32api.GetSystemMetrics(1)  # 获得屏幕分辨率Y轴
        print(sX)
        print(sY)

        """获取真实的分辨率"""
        hDC = win32gui.GetDC(0)
        w = win32print.GetDeviceCaps(hDC, win32con.DESKTOPHORZRES)  # 横向分辨率
        h = win32print.GetDeviceCaps(hDC, win32con.DESKTOPVERTRES)  # 纵向分辨率
        print(w)
        print(h)

        # 缩放比率
        screen_scale_rate = round(w / sX, 2)
        print(screen_scale_rate)
        print('windows')


def sum_lll(alist):  # 列总和
    l = len(alist[0])
    sumlist = []
    for i in range(l):
        suml = 0
        for j in range(len(alist)):
            suml += alist[j][i]
        sumlist.append(suml)
    return sumlist

class share_variable():
    def __init__(self):
        assert 'PYSHARE' in os.environ,'PYSHARE not in environment'
    def set(self,string):
        os.environ['PYSHARE'] = string
    def js_set(self,dict_):
        os.environ['PYSHARE'] = json.dumps(dict_)
    def __str__(self):
        return os.environ['PYSHARE']
    def __repr__(self):
        return os.environ['PYSHARE']
    def __int__(self):
        return int(self.__str__())
    def __float__(self):
        return float(self.__str__())
    def js(self):
        return json.loads(self.__str__())
    def __del__(self):
        os.environ['PYSHARE'] = 'NULL'


def transparent_png(path, xs, ys, new_name='transparent'):
    img = cv2_imread(path)
    h, w, d = img.shape

    if img.shape[2] != 4:
        png = np.ones((h, w, 1))*255
        # print('input jpg  take a few time to png')
        # png=np.zeros((h,w,4))
        # for n,i in enumerate(img):
        #     for m,j in enumerate(i):
        #         png[n][m]=np.matrix(img[n][m].tolist()+[255])
        png = np.append(img, png, axis=2)

    # print('pic shape: ',png.shape)

    png[ys[0]:ys[1], xs[0]:xs[1]] = np.matrix([0, 0, 0, 0])
    # print(os.path.join(path2dirname(path),new_name+'.png'))
    cv2.imwrite(os.path.join(path2dirname(path), new_name+'.png'), png)
    print('saved at', os.path.join(path2dirname(path), new_name+'.png'))


def truepath(file, *arg):
    return os.path.join(os.path.abspath(os.path.split(file)[0]), *arg)


def tensor2img(tensor, res=None):
    import torch
    if type(tensor) == type(torch.zeros(1)):
        try:
            n, a, b, c = tensor.shape

        except:
            a, b, c = tensor.shape
        if a < b:
            tensor = tensor.reshape(b, c, a)
        if tensor.mean() > 1:
            tensor = np.array(tensor.detach(), dtype=np.uint8)
        else:
            tensor = np.array(tensor.detach())
    if not res:
        try:
            cv2.imshow('n', tensor)
        except:
            cv2.imshow('w', tensor[:, :, 0])
    else:
        try:
            cv2.imshow('n', cv2.resize(tensor, res))
        except:
            cv2.imshow('w', cv2.resize(tensor[:, :, 0], res))
    cv2.waitKey(0)
    cv2.destroyAllWindows()


def write_file_name(file_dir,
                    output_place='this',
                    output_name='output.txt',
                    remove='.txt',
                    judge=None,
                    orpath=0):
    # 路径下文件的路径写入txt
    dir_list1 = os.listdir(file_dir)  # get_son
    # print(dir_list1)
    dir_list = []
    for name in dir_list1:
        if remove in name:
            dir_list1.remove(name)
            break
        if judge:
            if judge in name:
                dir_list.append(name)
        else:
            dir_list = dir_list1
    if output_place == 'this':
        p = ''
    elif output_place == 'that':
        p = str(file_dir) + '\\'
    fp = open(p + output_name, 'w+')
    # print(dir_list)
    for dir_name in dir_list:

        if dir_name == dir_list[-1]:
            if orpath:
                fp.write(file_dir + '\\' + str(dir_name))
            else:
                fp.write(str(dir_name))
        else:
            if orpath:
                fp.write('\\' + str(dir_name) + '\n')
            else:
                fp.write(str(dir_name) + '\n')
    fp.close()
    # print('write',len(dir_list),'items')


def zipath(zipfile_path, Date_file):
    import shutil
    if zipfile_path.endswith('.zip'):
        zipfile_path = zipfile_path.rstrip('.zip')

    shutil.make_archive(zipfile_path, 'zip', Date_file)
    print('zipfile saved at', zipfile_path)


def zipsort(xl, yl):
    z = list(zip(xl, yl))
    z.sort(key=lambda x: x[0])
    xe = [z[i][0] for i in range(len(z))]
    ye = [z[i][1] for i in range(len(z))]
    return xe, ye


######################################################################################################


def timety_(func):
    @wraps(func)
    def wrapper(*arg, **kw):
        t = time.time()
        r = func(*arg, **kw)
        print('function {} take {}s'.format(func.__name__, time.time()-t))
        return r
    return wrapper

def async_timety_(func):
    @wraps(func)
    async def wrapper(*arg, **kw):
        t = time.time()
        r = await func(*arg, **kw)
        print('function {} take {}s'.format(func.__name__, time.time()-t))
        return r
    return wrapper



######################################################################################################

if os.name=='nt':
    pysl_lib = [r'C:\ProgramData\Anaconda\Lib\site-packages',r'C:\ProgramData\Anaconda\envs\vtuber\Lib\site-packages'][0]
else:
    pysl_lib = '/usr/lib/python3/dist-packages/'
    
    
if __name__ == 'pysl':
    pysl_including = True
    current_time = getime()
    cfmt_print('@p@当前目录: {os.getcwd()} \n@e@当前时间: {current_time}')

    cfmt_print('@g@Module PYSL Importing')
    cfmt_print('@_@______________________________\n')

    try:
        assert get_file_same(os.path.join(pysl_lib, 'pysl.py'), __file__)
    except:
        cfmt_print(
           f'@^Rb@Module pysl may not update for pysl_lib ({pysl_lib})', return_=False)

elif __name__ == '__main__':
    current_time = getime()
    cfmt_print('@p@当前目录: {os.getcwd()} \n@e@当前时间: {current_time}')
    cfmt_print('@g@Module PYSL Already')
    cfmt_print('@_@______________________________\n')

    try:
        assert get_file_same(os.path.join(pysl_lib, 'pysl.py'), './pysl.py')
    except:
        if os.name == 'nt':
            cmd(f'copy {__file__} {pysl_lib}')
        else:
            cmd(f'cp {__file__} {pysl_lib}')
        cfmt_print('@^Hb@Update to pysl module')



######################################################################################################

# '{: >6}.jpg'.format(n)
# print(sys._getframe().f_code.co_name)
# sys._getframe().f_lineno
# .transpose(1,2,0))
# %* %cd% %user% %o %~dp0
# transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]
# cv2.CAP_DSHOW)
# __matmul__ @
# chcp 65001
# '\033[3A'
# .g 有效数字
# cv2 asarray bin dtype uint8
# ^([a-z]{3,}).*\1    ^(?:(?!Andrea).)*$     ^(?:(?!(?:^.?$|^(xx+?)\1+$)).)*$
# os.path.join(os.path.abspath(os.path.split(__file__)[0]),path,name)


# plt.plot(label='图例1')  plt.xlabel('坐标轴名') plt.bar(h) 柱状图
# plt.subplot(411 412 411 422)  plt.xticks(x,[xp1,xp2,xp3 ] plt.xlim(-1,9)范围
# plt.axis("on") plt.legend(loc=2)
# plt.title('ti',fontsize = 25) plt.savefig('wsnd')

# import sys
# t = sys.getfilesystemencoding()
# str.encode(xx).decode(t)
# import io
# sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')

# from io import BytesIO
# import gzip

# buff = BytesIO(rp_body)
# f = gzip.GzipFile(fileobj=buff)
# htmls = f.read().decode('utf-8')


# old=sys.stdout
# sys.stdout = codecs.lookup('iso8859-1')[-1]( sys.stdout)

# ensure_ascii=False,indent=2


# 		p2 v p1
# 		五子
# 		colorsboard touch to color
#       ftp 远程cmd
